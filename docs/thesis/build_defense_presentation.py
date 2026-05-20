from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn as docx_qn
from docx.shared import Cm as DocxCm, Pt as DocxPt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DIAGRAMS = ROOT / "diagrams"
OUT_PPTX = ROOT / "MindHelper_defense_presentation.pptx"
OUT_SCRIPT = ROOT / "MindHelper_defense_speech.docx"

WIDE_W = 13.333
WIDE_H = 7.5

NAVY = RGBColor(12, 25, 45)
BLUE = RGBColor(22, 116, 200)
CYAN = RGBColor(32, 178, 214)
INK = RGBColor(27, 39, 57)
MUTED = RGBColor(92, 108, 128)
LIGHT = RGBColor(244, 248, 252)
PALE = RGBColor(228, 241, 250)
GREEN = RGBColor(22, 153, 117)
ORANGE = RGBColor(235, 143, 55)
RED = RGBColor(205, 74, 74)
WHITE = RGBColor(255, 255, 255)


def Cm(value: float):
    """The slide design grid is written in inches; keep the short Cm() name for compact layout code."""
    return Inches(value)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[:2], 16), int(hex_value[2:4], 16), int(hex_value[4:], 16))


def set_run(run, size=18, color=INK, bold=False, font="Aptos"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, text, x, y, w, h, *, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        set_run(run, size=size, color=color, bold=bold, font=font)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, 0.9, 0.45, 11.8, 0.8, size=28, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, 0.95, 1.18, 11.0, 0.45, size=12, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0.92), Cm(1.72), Cm(1.3), Cm(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()


def add_footer(slide, number):
    add_textbox(slide, f"{number:02d}", 12.15, 6.95, 0.7, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    add_textbox(slide, "MindHelper", 0.9, 6.95, 2.0, 0.25, size=9, color=MUTED)


def add_bullets(slide, bullets, x, y, w, h, *, size=17, color=INK):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = color
    return box


def add_panel(slide, x, y, w, h, fill=WHITE, line=PALE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    panel = slide.shapes.add_shape(shape_type, Cm(x), Cm(y), Cm(w), Cm(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = line
    panel.line.width = Pt(1)
    return panel


def add_pill(slide, text, x, y, w, color, *, text_color=WHITE):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        set_run(run, size=10, color=text_color, bold=True)
    return pill


def add_image_fit(slide, path, x, y, w, h):
    path = str(path)
    pic = slide.shapes.add_picture(path, Cm(x), Cm(y), width=Cm(w))
    if pic.height > Cm(h):
        scale = Cm(h) / pic.height
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
    pic.left = Cm(x) + int((Cm(w) - pic.width) / 2)
    pic.top = Cm(y) + int((Cm(h) - pic.height) / 2)
    return pic


def add_background(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(13.333), Cm(0.12))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BLUE
    stripe.line.fill.background()


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    for i, color in enumerate([BLUE, CYAN, GREEN]):
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(8.0 + i * 0.9), Cm(0.8 + i * 1.15), Cm(4.3), Cm(4.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.fill.transparency = 22 + i * 12
        circle.line.fill.background()
    add_textbox(slide, "MindHelper", 0.9, 0.9, 4.0, 0.5, size=18, color=CYAN, bold=True)
    add_textbox(slide, "Вопросно-ответная система\nдля предварительной психологической диагностики\nна основе нейросети", 0.9, 1.65, 8.0, 2.2, size=28, color=WHITE, bold=True)
    add_textbox(slide, "Выпускная работа бакалавра\n09.03.02 Информационные системы и технологии", 0.95, 4.25, 6.2, 0.7, size=14, color=rgb("C8D7E8"))
    add_textbox(slide, "Кретов Даниил Владимирович\nКафедра программирования и информационных технологий\nНаучный руководитель: Шишко Ю.В., ассистент кафедры", 0.95, 5.55, 7.6, 0.85, size=11, color=rgb("E8F0F8"))
    add_image_fit(slide, DIAGRAMS / "03_safety_flow.png", 9.05, 4.65, 3.0, 1.65)
    add_textbox(slide, "Воронеж, 2026", 10.55, 6.75, 1.8, 0.3, size=10, color=rgb("C8D7E8"), align=PP_ALIGN.RIGHT)
    return slide


def add_problem(slide):
    add_title(slide, "Актуальность и проблема", "Психологическая поддержка требует доступности, но нейросеть должна быть ограничена безопасностью")
    add_bullets(slide, [
        "Пользователь может не обращаться к специалисту сразу и нуждается в первичном безопасном контакте.",
        "LLM умеет вести естественный диалог, но может неверно распознать риск или дать неподходящий совет.",
        "Нужна система, где модель помогает, но не принимает критические решения без safety-контура.",
    ], 0.95, 2.15, 6.2, 2.4)
    # Visual: tension scale
    add_panel(slide, 7.45, 2.05, 4.8, 2.8, fill=WHITE)
    add_textbox(slide, "Ключевое противоречие", 7.8, 2.35, 4.1, 0.35, size=15, color=NAVY, bold=True)
    add_pill(slide, "естественный диалог", 7.85, 3.0, 2.25, BLUE)
    add_pill(slide, "контроль риска", 9.95, 4.0, 1.9, RED)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Cm(8.25), Cm(3.55), Cm(3.0), Cm(0.45))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ORANGE
    arrow.line.fill.background()
    add_textbox(slide, "решение: гибридная архитектура", 8.05, 4.55, 3.9, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)


def add_goal(slide):
    add_title(slide, "Цель и задачи работы", "Сервис первичной поддержки: веб-интерфейс, Telegram-бот, база данных, LLM и safety-flow")
    add_panel(slide, 0.95, 2.0, 5.5, 3.25, fill=WHITE)
    add_textbox(slide, "Цель", 1.25, 2.28, 1.4, 0.35, size=17, color=BLUE, bold=True)
    add_textbox(slide, "Разработать вопросно-ответную систему, которая повышает доступность первичной психологической поддержки и снижает риск вредных ответов нейросети.", 1.25, 2.85, 4.8, 1.55, size=17, color=INK)
    add_panel(slide, 6.85, 2.0, 5.55, 3.25, fill=WHITE)
    add_bullets(slide, [
        "спроектировать ER-модель и PostgreSQL-базу;",
        "реализовать backend, frontend и Telegram-бот;",
        "интегрировать Qwen3 через Ollama;",
        "построить safety-flow и провести тестирование.",
    ], 7.15, 2.33, 5.0, 2.55, size=15)
    add_image_fit(slide, DIAGRAMS / "02_message_processing.png", 3.1, 5.35, 7.2, 0.95)


def add_model_choice(slide):
    add_title(slide, "Выбор нейросетевой модели", "В ходе исследования Qwen3 выбрана как практичный вариант для локального запуска")
    headers = ["Критерий", "ChatGPT", "Qwen3", "Llama", "Gemma"]
    rows = [
        ["Диалог", "высокий", "высокий", "средний", "средний"],
        ["Контроль данных", "внешний API", "локально", "локально", "локально"],
        ["Русский язык", "высокий", "хороший", "средний", "средний"],
        ["Интеграция", "зависимость", "Ollama", "сложнее", "ограниченно"],
        ["Итог", "не локально", "выбрана", "резерв", "резерв"],
    ]
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Cm(0.85), Cm(2.0), Cm(11.8), Cm(3.4)).table
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    for r in range(len(rows) + 1):
        for c in range(len(headers)):
            cell = table.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALE if r == 0 else WHITE
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run(run, size=10.5, color=NAVY if r == 0 else INK, bold=(r == 0 or (c == 2 and r == 5)))
    add_pill(slide, "Qwen3 + Ollama", 5.25, 5.85, 2.8, GREEN)


def add_architecture(slide):
    add_title(slide, "Общая архитектура сервиса", "Веб-сайт и Telegram-бот используют единый backend и общую базу данных")
    add_image_fit(slide, DIAGRAMS / "01_system_architecture.png", 0.75, 1.85, 7.25, 4.65)
    add_bullets(slide, [
        "единый chat-сервис для сайта и Telegram;",
        "PostgreSQL хранит пользователей, сообщения, события риска и справочники;",
        "нейросеть вызывается только после маршрутизации безопасности;",
        "администратор управляет моделью и актуальной информацией.",
    ], 8.25, 2.0, 4.15, 3.8, size=15)


def add_data(slide):
    add_title(slide, "Модель данных", "База данных обеспечивает историю диалога, аудит и управляемость справочников")
    add_image_fit(slide, DIAGRAMS / "08_er_model_large.png", 0.65, 1.75, 7.8, 4.85)
    add_panel(slide, 8.65, 1.95, 3.55, 3.7, fill=WHITE)
    add_bullets(slide, [
        "user_chat — один чат пользователя;",
        "chat_message — полная история диалога;",
        "crisis_event — события риска;",
        "emergency_resource — проверенные контакты помощи;",
        "neural_model_version — версия модели.",
    ], 8.95, 2.25, 3.0, 3.1, size=13.2)


def add_qwen(slide):
    add_title(slide, "Что находится под капотом Qwen3", "Модель генерирует текст по токенам, но safety-решение принимает не она")
    add_image_fit(slide, DIAGRAMS / "04_qwen_under_the_hood.png", 0.85, 1.85, 6.6, 4.65)
    add_panel(slide, 7.9, 2.0, 4.35, 3.6, fill=WHITE)
    add_bullets(slide, [
        "decoder-only Transformer;",
        "токенизация текста;",
        "контекст диалога;",
        "system prompt и policy layer;",
        "локальный запуск через Ollama.",
    ], 8.25, 2.35, 3.7, 2.75, size=15)


def add_safety(slide):
    add_title(slide, "Safety-flow", "Критический риск обрабатывается без свободной генерации модели")
    add_image_fit(slide, DIAGRAMS / "03_safety_flow.png", 0.8, 1.85, 7.0, 4.75)
    labels = [("low", GREEN), ("elevated", ORANGE), ("high", BLUE), ("critical", RED)]
    for i, (label, color) in enumerate(labels):
        add_pill(slide, label, 8.3 + (i % 2) * 1.95, 2.05 + (i // 2) * 0.8, 1.65, color)
    add_bullets(slide, [
        "low/elevated: практический совет + ограничения;",
        "high: уточняющий ASQ-сценарий;",
        "critical: шаблонный ответ и экстренные ресурсы из БД.",
    ], 8.2, 3.65, 4.0, 2.0, size=14.5)


def add_pipeline(slide):
    add_title(slide, "Формирование ответа", "Ответ модели проходит prompt-policy и response policy")
    add_image_fit(slide, DIAGRAMS / "05_response_pipeline.png", 0.8, 1.9, 7.0, 4.5)
    add_panel(slide, 8.15, 1.95, 4.1, 3.8, fill=WHITE)
    add_bullets(slide, [
        "модель не ставит диагноз;",
        "не назначает лекарства;",
        "не выдумывает контакты помощи;",
        "не отвечает свободно в critical-сценарии;",
        "сохраняет аудит маршрута.",
    ], 8.5, 2.25, 3.45, 3.1, size=14)


def add_implementation(slide):
    add_title(slide, "Программная реализация", "MindHelper реализован как полноценный сервис, а не отдельный чат-бот")
    lanes = [
        ("Backend", "Django\nDRF\nсервисный слой", BLUE),
        ("Frontend", "React\nчат\nкаталог специалистов", CYAN),
        ("Telegram", "бот\nединый chat API\nэкстренные ресурсы", GREEN),
        ("Admin", "версии модели\nсправочники\nsafety-аудит", ORANGE),
    ]
    for i, (title, body, color) in enumerate(lanes):
        x = 0.9 + i * 3.05
        add_panel(slide, x, 2.05, 2.55, 3.35, fill=WHITE)
        add_pill(slide, title, x + 0.35, 2.35, 1.85, color)
        add_textbox(slide, body, x + 0.25, 3.05, 2.0, 1.55, size=15, color=INK, align=PP_ALIGN.CENTER)
    add_image_fit(slide, DIAGRAMS / "01_system_architecture.png", 0.95, 5.55, 2.35, 0.85)
    add_textbox(slide, "Все каналы используют одну базу данных и единый safety-flow", 2.0, 5.85, 9.2, 0.45, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def add_testing(slide):
    add_title(slide, "Тестирование", "Проверялись функции сервиса и опасные сценарии")
    # Native bar-like visual with shapes
    items = [("Функциональные тесты", 82, BLUE), ("Safety-сценарии", 92, RED), ("Интеграция бота", 76, GREEN), ("Админка и БД", 80, ORANGE)]
    for i, (label, val, color) in enumerate(items):
        y = 2.05 + i * 0.82
        add_textbox(slide, label, 0.95, y, 3.0, 0.3, size=13.5, color=INK)
        base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(4.0), Cm(y + 0.04), Cm(5.4), Cm(0.22))
        base.fill.solid()
        base.fill.fore_color.rgb = PALE
        base.line.fill.background()
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(4.0), Cm(y + 0.04), Cm(5.4 * val / 100), Cm(0.22))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, f"{val}%", 9.65, y - 0.03, 0.8, 0.3, size=12, color=MUTED)
    add_panel(slide, 10.65, 1.85, 1.55, 3.7, fill=NAVY, line=NAVY)
    add_textbox(slide, "123", 10.83, 2.45, 1.15, 0.75, size=29, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "теста\nпройдено", 10.85, 3.18, 1.1, 0.75, size=11, color=rgb("DCE9F6"), align=PP_ALIGN.CENTER)
    add_image_fit(slide, DIAGRAMS / "03_safety_flow.png", 8.35, 5.0, 3.0, 0.75)
    add_textbox(slide, "Главный критерий: не пропустить critical-сценарий", 1.0, 5.7, 10.7, 0.45, size=17, color=NAVY, bold=True)


def add_results(slide):
    add_title(slide, "Результаты работы", "Получен прототип сервиса MindHelper и исследовательский safety-контур")
    results = [
        ("Сервис", "веб-интерфейс, Telegram-бот, backend API"),
        ("Данные", "PostgreSQL, ER-модель, справочники помощи"),
        ("Нейросеть", "Qwen3 через Ollama, версия модели в БД"),
        ("Безопасность", "risk levels, ASQ-сценарий, audit log"),
        ("Тестирование", "функциональные и red-team сценарии"),
    ]
    for i, (title, body) in enumerate(results):
        y = 1.9 + i * 0.72
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(1.05), Cm(y + 0.05), Cm(0.22), Cm(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CYAN if i % 2 == 0 else BLUE
        dot.line.fill.background()
        add_textbox(slide, title, 1.45, y - 0.02, 2.0, 0.3, size=14.5, color=NAVY, bold=True)
        add_textbox(slide, body, 3.2, y - 0.02, 8.2, 0.3, size=14, color=INK)
    add_panel(slide, 8.25, 5.35, 3.6, 0.72, fill=PALE)
    add_textbox(slide, "Основа для дальнейшей экспертной разметки и LoRA-адаптации", 8.45, 5.55, 3.2, 0.28, size=10.5, color=NAVY, align=PP_ALIGN.CENTER)
    add_image_fit(slide, DIAGRAMS / "05_response_pipeline.png", 8.25, 1.8, 3.65, 1.7)


def add_final(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_textbox(slide, "Спасибо за внимание", 0.95, 2.35, 8.2, 0.8, size=34, color=WHITE, bold=True)
    add_textbox(slide, "MindHelper: вопросно-ответная система предварительной психологической диагностики на основе нейросети", 1.0, 3.45, 8.3, 0.65, size=16, color=rgb("DCE9F6"))
    add_textbox(slide, "Кретов Даниил Владимирович\n09.03.02 Информационные системы и технологии", 1.0, 5.55, 5.7, 0.65, size=12, color=rgb("C8D7E8"))
    check = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(9.4), Cm(1.7), Cm(2.8), Cm(2.8))
    check.fill.solid()
    check.fill.fore_color.rgb = GREEN
    check.line.fill.background()
    add_textbox(slide, "Q&A", 9.83, 2.63, 1.9, 0.55, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Cm(WIDE_W * 2.54)
    prs.slide_height = Cm(WIDE_H * 2.54)
    # PowerPoint wide dimensions in cm above: 33.866 x 19.05, but python-pptx's Cm uses actual cm.
    # Reset to a standard 16:9 canvas explicitly.
    prs.slide_width = 12192000
    prs.slide_height = 6858000

    builders = [
        add_cover,
        lambda prs: build_content_slide(prs, add_problem, 2),
        lambda prs: build_content_slide(prs, add_goal, 3),
        lambda prs: build_content_slide(prs, add_model_choice, 4),
        lambda prs: build_content_slide(prs, add_architecture, 5),
        lambda prs: build_content_slide(prs, add_data, 6),
        lambda prs: build_content_slide(prs, add_qwen, 7),
        lambda prs: build_content_slide(prs, add_safety, 8),
        lambda prs: build_content_slide(prs, add_pipeline, 9),
        lambda prs: build_content_slide(prs, add_implementation, 10),
        lambda prs: build_content_slide(prs, add_testing, 11),
        lambda prs: build_content_slide(prs, add_results, 12),
        lambda prs: build_content_slide(prs, add_final, 13, footer=False),
    ]
    for builder in builders:
        builder(prs)
    prs.save(OUT_PPTX)


def build_content_slide(prs, content_builder, number, footer=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if content_builder is not add_final:
        add_background(slide)
    content_builder(slide)
    if footer:
        add_footer(slide, number)
    return slide


SCRIPT = [
    ("Слайд 1", "Уважаемая Государственная экзаменационная комиссия. Позвольте представить Вашему вниманию выпускную работу бакалавра по направлению 09.03.02 Информационные системы и технологии. Я, Кретов Даниил Владимирович, студент кафедры программирования и информационных технологий. Научный руководитель — Шишко Юлия Владимировна, ассистент кафедры программирования и информационных технологий. Тема работы — «Разработка вопросно-ответной системы для предварительной психологической диагностики на основе нейросети»."),
    ("Слайд 2", "Актуальность работы связана с потребностью в доступном первом канале психологической поддержки. Пользователь может находиться в тревожном или кризисном состоянии и не сразу обратиться к специалисту. Большие языковые модели способны поддерживать естественный диалог, но без ограничений могут неверно распознать риск или дать неподходящий ответ. Поэтому задача работы состоит не только в создании чат-бота, а в построении безопасной системы."),
    ("Слайд 3", "Цель работы — разработать вопросно-ответную систему, которая помогает пользователю описать состояние, получить первичную безопасную обратную связь и сведения о ресурсах помощи. Для достижения цели были поставлены задачи: спроектировать базу данных, реализовать backend и frontend, добавить Telegram-бота, интегрировать локальную нейросетевую модель и построить safety-flow для ограничения опасных ответов."),
    ("Слайд 4", "В ходе исследования рассматривались разные варианты языковых моделей. ChatGPT показывает высокое качество диалога, но требует внешнего API и передачи данных провайдеру. Для дипломного проекта важнее локальный запуск и контроль данных, поэтому была выбрана Qwen3 через Ollama. Она дала достаточное качество русскоязычного диалога и позволила встроить модель в собственный backend."),
    ("Слайд 5", "На этом слайде представлена общая архитектура сервиса. Пользователь может работать через web-интерфейс или Telegram-бота. Оба канала обращаются к одному backend. Backend сохраняет данные в PostgreSQL, вызывает safety-flow, при допустимом маршруте обращается к Ollama и сохраняет результат. Администратор управляет справочниками, экстренными ресурсами и версиями модели."),
    ("Слайд 6", "На этом слайде показана модель данных. История диалога хранится в сущности chat_message, а не только в оперативной памяти приложения. Отдельно сохраняются кризисные события, экстренные ресурсы, специалисты, версии нейросети и записи safety-аудита. Такая структура позволяет анализировать работу системы, обновлять справочники и документировать поведение модели."),
    ("Слайд 7", "Qwen3 используется как языковой компонент, но не как единственный центр принятия решения. Модель работает с токенами и контекстом диалога, формируя следующий текстовый ответ. Однако перед обращением к модели сообщение проходит программную оценку риска. Это важно, потому что в психологически чувствительном диалоге нельзя полностью доверять свободной генерации модели."),
    ("Слайд 8", "На этом слайде представлен safety-flow. Система выделяет четыре уровня риска: low, elevated, high и critical. В обычных сценариях модель может дать практический поддерживающий совет. При высоком риске запускается уточняющий сценарий. В critical-сценарии свободная генерация блокируется, а пользователь получает заранее подготовленный ответ и экстренные контакты из базы данных."),
    ("Слайд 9", "Ответ нейросети формируется через несколько ограничений. Сначала задается system prompt, затем выбирается сценарий состояния, после этого выполняется запрос к Qwen3 через Ollama. Готовый ответ дополнительно проходит response policy. Модель не должна ставить диагноз, назначать лекарства, выдумывать контакты помощи или давать вредные инструкции."),
    ("Слайд 10", "На этом слайде показаны основные компоненты реализации. Backend разработан на Django и Django REST Framework. Frontend реализует пользовательский сайт с чатом и каталогом. Telegram-бот служит дополнительным каналом доступа. Django Admin используется для управления версиями модели, экстренными ресурсами и журналом safety-аудита. Все компоненты работают через единый сервисный слой."),
    ("Слайд 11", "Тестирование включало функциональные и safety-сценарии. Проверялись регистрация, уникальность email, API сообщений, сохранение истории, работа Telegram-бота, справочники и административные функции. Отдельно проверялись опасные формулировки, при которых система должна выбрать high или critical-маршрут. Главный критерий безопасности — не пропустить critical-сценарий."),
    ("Слайд 12", "В результате был получен прототип сервиса MindHelper. Реализованы web-интерфейс, Telegram-бот, backend API, PostgreSQL-база данных, интеграция Qwen3 через Ollama, safety-flow, red-team сценарии, журнал аудита и административное управление. Система может быть развита дальше за счет экспертной разметки обезличенных диалогов и последующего улучшения модели."),
    ("Слайд 13", "Спасибо за внимание."),
]


def build_script_docx():
    doc = Document()
    section = doc.sections[0]
    section.left_margin = DocxCm(3)
    section.right_margin = DocxCm(1.5)
    section.top_margin = DocxCm(2)
    section.bottom_margin = DocxCm(2)
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Текст выступления к презентации")
    run.font.name = "Times New Roman"
    run._element.rPr.rFonts.set(docx_qn("w:eastAsia"), "Times New Roman")
    run.font.size = DocxPt(14)
    run.bold = True
    for slide, text in SCRIPT:
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.5
        p.paragraph_format.first_line_indent = Cm(0)
        r = p.add_run(f"{slide}: ")
        r.bold = True
        r.font.name = "Times New Roman"
        r._element.rPr.rFonts.set(docx_qn("w:eastAsia"), "Times New Roman")
        r.font.size = DocxPt(14)
        r = p.add_run(text)
        r.font.name = "Times New Roman"
        r._element.rPr.rFonts.set(docx_qn("w:eastAsia"), "Times New Roman")
        r.font.size = DocxPt(14)
    doc.save(OUT_SCRIPT)


if __name__ == "__main__":
    build_presentation()
    build_script_docx()
    print(OUT_PPTX)
    print(OUT_SCRIPT)
