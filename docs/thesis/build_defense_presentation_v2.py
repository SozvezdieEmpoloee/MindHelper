# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt as PptPt

ROOT = Path(__file__).resolve().parent
OUT_PPTX = ROOT / "MindHelper_defense_presentation_v2.pptx"
OUT_SCRIPT = ROOT / "MindHelper_defense_speech_v2.docx"

W, H = 13.333, 7.5
NAVY = RGBColor(12, 25, 45)
INK = RGBColor(30, 42, 62)
MUTED = RGBColor(91, 108, 128)
BLUE = RGBColor(20, 116, 205)
CYAN = RGBColor(32, 178, 214)
GREEN = RGBColor(22, 153, 117)
ORANGE = RGBColor(235, 143, 55)
RED = RGBColor(205, 74, 74)
LIGHT = RGBColor(246, 250, 253)
PALE = RGBColor(226, 239, 249)
WHITE = RGBColor(255, 255, 255)
SOFT_GREEN = RGBColor(222, 245, 238)
SOFT_RED = RGBColor(249, 225, 225)
SOFT_ORANGE = RGBColor(255, 238, 218)


def X(v): return Inches(v)


def set_text(run, size=16, color=INK, bold=False, font="Aptos"):
    run.font.name = font
    run.font.size = PptPt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(X(x), X(y), X(w), X(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = X(0.04)
    tf.margin_right = X(0.04)
    tf.margin_top = X(0.02)
    tf.margin_bottom = X(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for r in p.runs:
        set_text(r, size, color, bold, font)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.75, 0.42, 11.9, 0.55, 27, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.78, 1.08, 11.7, 0.38, 12.5, MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, X(0.78), X(1.55), X(1.35), X(0.055))
    line.fill.solid(); line.fill.fore_color.rgb = CYAN; line.line.fill.background()


def add_footer(slide, n):
    add_text(slide, "MindHelper", 0.78, 7.05, 1.8, 0.23, 8.5, MUTED)
    add_text(slide, f"{n:02d}", 12.18, 7.05, 0.55, 0.23, 8.5, MUTED, align=PP_ALIGN.RIGHT)


def bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, X(0), X(0), X(W), X(0.10))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = BLUE; stripe.line.fill.background()


def panel(slide, x, y, w, h, fill=WHITE, line=PALE, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, X(x), X(y), X(w), X(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = PptPt(1)
    return shp


def pill(slide, text, x, y, w, color, text_color=WHITE, size=10.5):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, X(x), X(y), X(w), X(0.42))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    for r in p.runs: set_text(r, size, text_color, True)
    return shp


def bullets(slide, items, x, y, w, h, size=14.5, color=INK):
    box = slide.shapes.add_textbox(X(x), X(y), X(w), X(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.space_after = PptPt(7)
        p.font.name = "Aptos"; p.font.size = PptPt(size); p.font.color.rgb = color
    return box


def arrow(slide, x1, y1, x2, y2, color=BLUE):
    conn = slide.shapes.add_connector(1, X(x1), X(y1), X(x2), X(y2))
    conn.line.color.rgb = color; conn.line.width = PptPt(2.0)
    return conn


def entity(slide, name, fields, x, y, w, h, accent=BLUE):
    panel(slide, x, y, w, h, WHITE, PALE)
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, X(x), X(y), X(w), X(0.46))
    head.fill.solid(); head.fill.fore_color.rgb = accent; head.line.fill.background()
    tf = head.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = name; p.alignment = PP_ALIGN.CENTER
    for r in p.runs: set_text(r, 11.5, WHITE, True)
    txt = "\n".join(fields)
    add_text(slide, txt, x+0.12, y+0.58, w-0.24, h-0.68, 9.7, INK)


def metric(slide, value, label, x, y, color=BLUE):
    add_text(slide, value, x, y, 1.45, 0.55, 25, color, True, PP_ALIGN.CENTER)
    add_text(slide, label, x-0.05, y+0.58, 1.55, 0.35, 9.5, MUTED, align=PP_ALIGN.CENTER)


def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, X(8.4), X(0), X(4.95), X(7.5))
    band.fill.solid(); band.fill.fore_color.rgb = RGBColor(18, 54, 86); band.line.fill.background()
    for i, c in enumerate([BLUE, CYAN, GREEN]):
        o = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, X(8.2+i*0.58), X(1.15+i*0.78), X(3.8), X(3.8))
        o.fill.solid(); o.fill.fore_color.rgb = c; o.fill.transparency = 26+i*10; o.line.fill.background()
    add_text(s, "MindHelper", 0.85, 0.82, 4.0, 0.4, 18, CYAN, True)
    add_text(s, "Вопросно-ответная система\nдля предварительной психологической диагностики\nна основе нейросети", 0.82, 1.55, 7.7, 2.25, 27, WHITE, True)
    add_text(s, "Выпускная работа бакалавра\n09.03.02 Информационные системы и технологии", 0.86, 4.22, 6.3, 0.65, 13.5, RGBColor(205, 220, 237))
    add_text(s, "Кретов Даниил Владимирович\nКафедра программирования и информационных технологий\nНаучный руководитель: Шишко Ю.В., ассистент кафедры", 0.86, 5.55, 7.5, 0.78, 10.8, RGBColor(232, 240, 248))
    pill(s, "LLM + safety-flow", 9.42, 4.95, 2.45, GREEN)
    add_text(s, "Воронеж, 2026", 10.7, 6.82, 1.7, 0.28, 10, RGBColor(205, 220, 237), align=PP_ALIGN.RIGHT)


def slide_problem(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Актуальность и проблема", "Цифровая поддержка должна быть доступной, но особенно осторожной в кризисных ситуациях")
    bullets(s, [
        "Пользователь может нуждаться в первичном безопасном контакте до обращения к специалисту.",
        "Большая языковая модель способна вести естественный диалог, но не является врачом и может ошибаться в рискованных случаях.",
        "Для психологической тематики важна не только генерация ответа, но и контролируемая маршрутизация опасных сценариев."
    ], 0.9, 2.05, 5.65, 2.5, 15.5)
    add_text(s, "Ключевая идея работы", 7.35, 2.0, 4.7, 0.35, 17, NAVY, True, PP_ALIGN.CENTER)
    pill(s, "естественный диалог", 7.35, 2.75, 2.15, BLUE)
    pill(s, "проверка риска", 10.1, 2.75, 1.95, ORANGE)
    ar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, X(8.7), X(3.48), X(2.1), X(0.48)); ar.fill.solid(); ar.fill.fore_color.rgb = CYAN; ar.line.fill.background()
    pill(s, "безопасная рекомендация", 8.02, 4.35, 3.35, GREEN)
    add_text(s, "Система строится как управляемый сервис, где нейросеть помогает в диалоге, а safety-контур ограничивает опасные ответы.", 7.2, 5.18, 4.85, 0.78, 13.2, MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, n)


def slide_goal(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Цель и задачи работы", "Цель уточнена как инженерная и исследовательская: помощь пользователю плюс контроль безопасности")
    panel(s, 0.85, 1.95, 11.65, 1.42, WHITE)
    add_text(s, "Цель", 1.12, 2.18, 1.0, 0.3, 16, BLUE, True)
    add_text(s, "Разработать прототип вопросно-ответной системы, повышающей доступность первичной психологической поддержки за счёт локально интегрированной нейросети и формализованного safety-контура, который снижает риск вредных рекомендаций и обеспечивает фиксацию критических сценариев.", 2.05, 2.03, 10.05, 0.95, 15, INK)
    tasks = [
        ("1", "Проанализировать предметную область и риски применения LLM в психологической поддержке."),
        ("2", "Спроектировать ER-модель, PostgreSQL-базу и хранение истории диалога."),
        ("3", "Реализовать веб-интерфейс, Django backend и Telegram-бот как единый chat-сервис."),
        ("4", "Интегрировать Qwen3 через Ollama и хранить сведения о версии модели в БД."),
        ("5", "Построить safety-flow: уровни риска, ASQ-сценарий, policy layer и аудит маршрутов."),
        ("6", "Провести функциональное, интеграционное и red-team тестирование опасных сценариев."),
    ]
    y = 3.85
    for i, (num, text) in enumerate(tasks):
        x = 0.95 if i < 3 else 6.75
        yy = y + (i % 3) * 0.82
        pill(s, num, x, yy, 0.42, BLUE, size=10)
        add_text(s, text, x+0.55, yy-0.03, 5.05, 0.55, 12.2, INK)
    add_footer(s, n)


def slide_model_choice(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Выбор нейросетевой модели", "Сравнение проводилось по применимости к локальному сервису психологической поддержки")
    headers = ["Модель", "Параметры / слои", "Контекст", "Локальный запуск", "Вывод для проекта"]
    rows = [
        ["Qwen3-8B", "8,2B / 36", "40 960 токенов", "Ollama, 4-bit; подходит для RTX 3060 Ti", "выбрана: баланс качества, русского диалога и управляемости"],
        ["Llama 3.1 8B", "8B / 32", "128K", "возможен, но лицензия и русская специализация менее удобны", "резервный вариант"],
        ["Gemma 2 9B", "9B / 42", "8K", "возможен, но контекст короче", "ограничение для длинной истории диалога"],
        ["Mistral 7B v0.3", "7B / 32", "32K", "лёгкий запуск, Apache 2.0", "хорош как быстрый baseline"],
    ]
    table = s.shapes.add_table(1+len(rows), len(headers), X(0.62), X(1.92), X(12.1), X(3.95)).table
    widths = [1.55, 1.75, 1.65, 3.25, 3.9]
    for i,wid in enumerate(widths): table.columns[i].width = X(wid)
    for c,h in enumerate(headers): table.cell(0,c).text = h
    for r,row in enumerate(rows,1):
        for c,val in enumerate(row): table.cell(r,c).text = val
    for r in range(len(rows)+1):
        for c in range(len(headers)):
            cell = table.cell(r,c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid(); cell.fill.fore_color.rgb = PALE if r==0 else (SOFT_GREEN if r==1 else WHITE)
            cell.margin_left = X(0.04); cell.margin_right = X(0.04); cell.margin_top = X(0.02); cell.margin_bottom = X(0.02)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c<3 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_text(run, 8.7 if r else 9.5, NAVY if r==0 else INK, r==0 or (r==1 and c==0))
    add_text(s, "Критерии выбора: архитектурная ёмкость, длина контекста, возможность локального запуска, качество русскоязычного диалога, интеграция с safety-контуром.", 0.8, 6.12, 11.7, 0.45, 12.2, MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, n)


def slide_architecture(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Общая архитектура сервиса", "Сайт и Telegram-бот работают через единый backend, поэтому логика безопасности не дублируется")
    labels = [
        ("Пользователь", 0.95, 2.25, BLUE), ("Web / React", 2.85, 1.55, CYAN), ("Telegram-бот", 2.85, 3.05, CYAN),
        ("Django REST API", 5.45, 2.25, BLUE), ("Chat service", 7.75, 1.55, GREEN), ("Safety router", 7.75, 3.05, ORANGE),
        ("Ollama + Qwen3", 10.25, 1.55, NAVY), ("PostgreSQL", 10.25, 3.05, NAVY)
    ]
    for text,x,y,c in labels:
        panel(s,x,y,1.85,0.72,WHITE,PALE); pill(s,text,x+0.08,y+0.16,1.69,c,size=9)
    arrow(s,2.25,2.61,2.85,1.9); arrow(s,2.25,2.61,2.85,3.4)
    arrow(s,4.7,1.9,5.45,2.55); arrow(s,4.7,3.4,5.45,2.55)
    arrow(s,7.3,2.55,7.75,1.9); arrow(s,7.3,2.55,7.75,3.4)
    arrow(s,9.6,1.9,10.25,1.9); arrow(s,9.6,3.4,10.25,3.4)
    add_text(s, "Принцип реализации", 0.95, 5.1, 2.25, 0.35, 16, NAVY, True)
    bullets(s, [
        "все сообщения проходят общий сервис обработки диалога;",
        "история и события риска сохраняются в базе данных;",
        "нейросеть вызывается только после проверки входного сообщения;",
        "критические сценарии не отдаются на свободную генерацию модели."
    ], 0.95, 5.55, 11.2, 0.9, 12.4)
    add_footer(s, n)


def slide_er_chat(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "ER-модель: пользователь, канал и диалог", "Диаграмма разбита на блоки, чтобы показать назначение сущностей без перегрузки")
    entity(s, "user_account", ["id UUID", "email", "password_hash", "role", "created_at"], 0.82, 2.0, 2.25, 2.0, BLUE)
    entity(s, "channel_account", ["id UUID", "user_id", "channel_type", "external_id", "display_name"], 3.75, 2.0, 2.25, 2.0, CYAN)
    entity(s, "user_chat", ["id UUID", "user_id", "status", "created_at", "updated_at"], 6.65, 2.0, 2.25, 2.0, GREEN)
    entity(s, "chat_message", ["id UUID", "chat_id", "sender_type", "content", "created_at"], 9.55, 2.0, 2.25, 2.0, ORANGE)
    arrow(s,3.08,3.0,3.75,3.0); arrow(s,6.0,3.0,6.65,3.0); arrow(s,8.9,3.0,9.55,3.0)
    entity(s, "neural_model_version", ["id UUID", "provider", "model_name", "version_tag", "is_active"], 4.1, 4.72, 3.0, 1.75, NAVY)
    add_text(s, "Каждый ответ ассистента может быть связан с актуальной версией модели. Это нужно для воспроизводимости тестов и анализа качества ответов.", 7.45, 5.0, 4.25, 1.0, 12.3, MUTED)
    add_footer(s, n)


def slide_er_safety(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "ER-модель: безопасность и справочники", "Отдельный блок хранит маршруты риска, проверенные ресурсы помощи и данные для интерфейса")
    entity(s, "crisis_event", ["id UUID", "chat_id", "message_id", "risk_level", "screening_status", "status"], 0.82, 1.95, 2.55, 2.15, RED)
    entity(s, "safety_route_audit", ["id UUID", "event_id", "route_code", "escalation_action", "human_review_flag"], 3.88, 1.95, 2.8, 2.15, ORANGE)
    entity(s, "emergency_resource", ["id UUID", "title", "phone", "region", "is_active"], 7.22, 1.95, 2.55, 2.15, GREEN)
    entity(s, "specialist_location", ["id UUID", "specialist_id", "address", "latitude", "longitude"], 10.25, 1.95, 2.55, 2.15, CYAN)
    arrow(s,3.37,3.0,3.88,3.0); arrow(s,6.68,3.0,7.22,3.0); arrow(s,9.77,3.0,10.25,3.0)
    panel(s, 1.1, 4.82, 10.95, 1.05, WHITE)
    add_text(s, "Зачем это нужно", 1.42, 5.05, 1.7, 0.3, 15, BLUE, True)
    add_text(s, "Свободный текст модели не используется как единственный источник решения: база хранит факт риска, маршрут обработки, действие эскалации и проверенные контакты помощи, которые показываются пользователю на сайте и в Telegram.", 3.1, 4.98, 8.55, 0.55, 12.5, INK)
    add_footer(s, n)


def slide_qwen(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Что находится под капотом Qwen3", "Нейросеть рассматривается как генератор текста, а не как автономный диагностический механизм")
    blocks = [
        ("Токенизация", "фраза пользователя превращается в последовательность токенов", BLUE),
        ("Transformer decoder", "36 слоёв обрабатывают контекст и строят вероятности следующего токена", CYAN),
        ("Контекст диалога", "модель получает ограниченную историю сообщений и системные инструкции", GREEN),
        ("Генерация", "ответ формируется последовательно, токен за токеном", ORANGE),
        ("Policy layer", "после генерации ответ проверяется и при необходимости заменяется безопасным шаблоном", RED),
    ]
    x = 0.78
    for title, body, color in blocks:
        panel(s, x, 2.25, 2.25, 2.35, WHITE)
        pill(s, title, x+0.18, 2.48, 1.89, color, size=9)
        add_text(s, body, x+0.18, 3.15, 1.88, 1.0, 10.4, INK, align=PP_ALIGN.CENTER)
        x += 2.5
    for x in [2.95, 5.45, 7.95, 10.45]:
        ar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, X(x), X(3.18), X(0.42), X(0.28))
        ar.fill.solid(); ar.fill.fore_color.rgb = MUTED; ar.line.fill.background()
    add_text(s, "В проекте Qwen3 запускается локально через Ollama. Это снижает зависимость от внешнего API и позволяет фиксировать версию модели в базе данных.", 1.1, 5.55, 11.2, 0.65, 13, MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, n)


def slide_safety(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Safety-flow", "Контур безопасности построен как обязательный слой до и после обращения к модели")
    levels = [("low", SOFT_GREEN, GREEN, "обычный поддерживающий диалог"), ("elevated", SOFT_ORANGE, ORANGE, "уточнение состояния и мягкая рекомендация"), ("high", SOFT_RED, RED, "ASQ-сценарий и показ ресурсов помощи"), ("critical", SOFT_RED, RED, "блок свободной генерации, экстренные контакты")]
    y = 2.0
    for i,(name,fill,c,desc) in enumerate(levels):
        x = 0.9 + i*3.0
        panel(s,x,y,2.55,1.55,fill,c)
        add_text(s,name,x+0.2,y+0.25,2.1,0.3,17,c,True,PP_ALIGN.CENTER)
        add_text(s,desc,x+0.22,y+0.78,2.1,0.45,10.5,INK,align=PP_ALIGN.CENTER)
    bullets(s, [
        "Входной текст проверяется правилами и сценариями риска.",
        "При high/critical риск фиксируется в crisis_event и safety_route_audit.",
        "Для critical-сценария модель не даёт свободный совет; пользователь получает экстренные ресурсы из БД.",
        "После ответа выполняется response policy: запрет диагнозов, лекарственных назначений и выдуманных контактов."
    ], 1.05, 4.25, 11.1, 1.35, 12.7)
    add_footer(s, n)


def slide_answer(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Формирование ответа", "Ответ не является сырым выводом модели: он проходит несколько контролируемых этапов")
    steps = [("1", "Сообщение", BLUE), ("2", "Оценка риска", ORANGE), ("3", "Prompt-policy", GREEN), ("4", "Qwen3", NAVY), ("5", "Response policy", RED), ("6", "Ответ пользователю", BLUE)]
    x=0.85
    for num,name,c in steps:
        panel(s,x,2.55,1.75,1.0,WHITE,PALE)
        pill(s,num,x+0.13,2.82,0.35,c,size=9)
        add_text(s,name,x+0.55,2.73,1.0,0.35,10.5,INK,True,PP_ALIGN.CENTER)
        x+=2.05
    for x in [2.58,4.63,6.68,8.73,10.78]:
        ar=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,X(x),X(2.9),X(0.35),X(0.22)); ar.fill.solid(); ar.fill.fore_color.rgb=MUTED; ar.line.fill.background()
    panel(s, 1.3, 4.35, 10.75, 1.25, WHITE)
    add_text(s, "Ограничения ответа", 1.65, 4.65, 2.1, 0.3, 15, RED, True)
    add_text(s, "Система не ставит диагноз, не назначает лекарства, не романтизирует самоповреждение, не выдумывает телефоны помощи и не оставляет critical-сценарий без маршрутизации.", 3.6, 4.55, 7.95, 0.5, 12.5, INK)
    add_footer(s, n)


def slide_implementation(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Программная реализация", "Нейросетевой контур показан как отдельная часть программной реализации")
    cards = [
        ("Backend", "Django REST Framework\nсервисный слой\nтранзакции БД", BLUE),
        ("Frontend", "React\nчат пользователя\nкарта специалистов", CYAN),
        ("Telegram", "отдельный polling-процесс\nобщий chat service\nресурсы помощи из БД", GREEN),
        ("Нейросеть", "Ollama + Qwen3\nmodel_version в БД\nprompt-policy\nresponse policy", ORANGE),
    ]
    for i,(t,b,c) in enumerate(cards):
        x=0.9+i*3.05
        panel(s,x,2.0,2.55,2.45,WHITE,PALE)
        pill(s,t,x+0.35,2.28,1.85,c,size=10.5)
        add_text(s,b,x+0.24,3.0,2.08,0.95,11.2,INK,align=PP_ALIGN.CENTER)
    add_text(s, "Нейросетевая часть включает локальный запуск Qwen3, хранение версии модели, сбор контекста диалога, prompt-policy и финальную проверку response policy перед показом ответа пользователю.", 1.1, 5.25, 11.1, 0.85, 12.3, MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, n)


def slide_testing(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s, "Тестирование", "Проверялась не только работоспособность интерфейса, но и устойчивость к вредным сценариям")
    metric(s,"123","автоматизированных теста",1.2,2.05,BLUE)
    metric(s,"4","уровня риска",3.5,2.05,ORANGE)
    metric(s,"0","допустимых пропусков critical",5.75,2.05,RED)
    metric(s,"2","канала общения",8.0,2.05,GREEN)
    panel(s, 1.0, 4.0, 11.25, 1.45, WHITE)
    bullets(s,[
        "unit-тесты сервисов chat и neural_engine;",
        "интеграционные тесты API, Telegram-бота и PostgreSQL;",
        "red-team набор фраз для проверки самоповреждения, угрозы жизни и опасных советов;",
        "проверка того, что critical-сценарии переводятся в emergency-response без свободной генерации."
    ],1.3,4.28,10.5,0.95,12.2)
    add_footer(s,n)


def slide_results(prs,n):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    add_title(s,"Результаты работы","Получен прототип сервиса и исследовательский контур безопасной нейросетевой консультации")
    items=[("Сервис","веб-интерфейс, Telegram-бот, Django API",BLUE), ("Данные","PostgreSQL, ER-модель, история диалога",CYAN), ("Нейросеть","Qwen3 через Ollama, версия модели в БД",ORANGE), ("Безопасность","risk levels, ASQ-сценарий, safety audit",RED), ("Тестирование","функциональные и red-team сценарии",GREEN)]
    for i,(t,b,c) in enumerate(items):
        y=1.95+i*0.85
        pill(s,t,1.1,y,1.65,c,size=9.5)
        add_text(s,b,3.0,y+0.03,7.8,0.35,13.2,INK)
    panel(s,1.1,6.05,10.9,0.65,WHITE)
    add_text(s,"Практический результат: основа полноценного сервиса первичной психологической поддержки с возможностью дальнейшей экспертной разметки и LoRA-адаптации.",1.4,6.2,10.3,0.3,11.8,MUTED,align=PP_ALIGN.CENTER)
    add_footer(s,n)


def slide_final(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb=NAVY
    add_text(s,"Спасибо за внимание",0.9,2.2,8.0,0.7,34,WHITE,True)
    add_text(s,"MindHelper: вопросно-ответная система предварительной психологической диагностики на основе нейросети",0.95,3.15,8.2,0.55,15,RGBColor(210,225,240))
    add_text(s,"Кретов Даниил Владимирович\n09.03.02 Информационные системы и технологии",0.95,5.4,5.3,0.55,12,RGBColor(210,225,240))
    pill(s,"Q&A",9.7,4.9,1.65,CYAN,size=16)


def build_deck():
    prs=Presentation(); prs.slide_width=X(W); prs.slide_height=X(H)
    cover(prs)
    funcs=[slide_problem,slide_goal,slide_model_choice,slide_architecture,slide_er_chat,slide_er_safety,slide_qwen,slide_safety,slide_answer,slide_implementation,slide_testing,slide_results]
    for idx,fn in enumerate(funcs,2): fn(prs,idx)
    slide_final(prs)
    prs.save(OUT_PPTX)


def build_speech():
    paras=[
        "Слайд 1: Уважаемая Государственная экзаменационная комиссия. Позвольте представить Вашему вниманию выпускную работу бакалавра по направлению 09.03.02 Информационные системы и технологии. Я, Кретов Даниил Владимирович, студент кафедры программирования и информационных технологий. Научный руководитель — Шишко Юлия Владимировна, ассистент кафедры программирования и информационных технологий. Тема работы — «Разработка вопросно-ответной системы для предварительной психологической диагностики на основе нейросети».",
        "Слайд 2: Актуальность работы связана с тем, что человеку может быть нужна первичная психологическая поддержка до обращения к специалисту. Большие языковые модели позволяют вести естественный диалог, однако в этой предметной области особенно важны ограничения: модель может ошибиться в оценке риска или дать неподходящий совет. Поэтому в работе решается задача не просто чат-бота, а управляемой системы с обязательным контуром безопасности.",
        "Слайд 3: Цель работы состоит в разработке прототипа вопросно-ответной системы, которая повышает доступность первичной психологической поддержки за счёт локально интегрированной нейросети и формализованного safety-контура. Для достижения цели были поставлены задачи анализа предметной области, проектирования базы данных, реализации сайта и Telegram-бота, интеграции Qwen3, построения safety-flow и тестирования опасных сценариев.",
        "Слайд 4: Для выбора нейросетевой модели были использованы не только общие бенчмарки, но и прикладные критерии проекта: количество параметров и слоёв, длина контекста, доступность локального запуска, качество русскоязычного диалога и удобство интеграции. По совокупности критериев была выбрана Qwen3-8B, так как она обеспечивает приемлемый баланс качества, локального запуска через Ollama и управляемости в составе safety-контура.",
        "Слайд 5: Архитектура сервиса построена вокруг единого backend. Веб-интерфейс и Telegram-бот обращаются к одному chat-сервису, поэтому история диалога, проверка риска и обращение к нейросети реализованы единообразно. Это снижает дублирование логики и позволяет применять один и тот же контур безопасности во всех каналах общения.",
        "Слайд 6: Первая часть ER-модели описывает пользователя, канал связи и диалог. Сущность user_account хранит учётную запись, channel_account связывает пользователя с внешним каналом, user_chat фиксирует чат, а chat_message хранит полную историю сообщений. Отдельно хранится neural_model_version, чтобы понимать, какая версия модели участвовала в формировании ответа.",
        "Слайд 7: Вторая часть ER-модели отвечает за безопасность и справочники. crisis_event фиксирует событие риска, safety_route_audit хранит маршрут обработки, emergency_resource содержит проверенные контакты помощи, а specialist_location используется для отображения специалистов и организаций на карте. Такое разделение делает базу данных читаемой и расширяемой.",
        "Слайд 8: Qwen3 в проекте рассматривается как генератор текста на основе decoder-only Transformer. Сообщение пользователя токенизируется, проходит через слои модели, после чего ответ формируется последовательно по токенам. При этом решение о безопасности не делегируется модели полностью: перед генерацией и после неё работают policy-слои.",
        "Слайд 9: Safety-flow включает четыре уровня риска: low, elevated, high и critical. Для обычных и умеренных сценариев система формирует поддерживающий ответ и практическую рекомендацию. Для high используется уточняющий сценарий, а для critical свободная генерация блокируется, создаётся событие риска и пользователю показываются экстренные ресурсы из базы данных.",
        "Слайд 10: Формирование ответа состоит из нескольких этапов. Сначала сообщение пользователя оценивается по риску, затем формируется системный prompt с ограничениями, после чего вызывается Qwen3. Полученный текст проходит response policy: система не допускает постановку диагноза, лекарственные назначения, выдуманные контакты помощи и опасные рекомендации.",
        "Слайд 11: Программная реализация включает frontend на React, backend на Django REST Framework, Telegram-бота как отдельный процесс и нейросетевой контур Ollama плюс Qwen3. На этом слайде отдельно выделена нейросетевая часть: версия модели хранится в базе данных, а её ответы проходят prompt-policy и response policy перед показом пользователю.",
        "Слайд 12: Тестирование включало функциональные и интеграционные проверки, а также red-team сценарии. Особое внимание уделялось критическим сообщениям: система должна распознавать опасные формулировки, фиксировать кризисное событие и не отдавать такие случаи на свободную генерацию модели.",
        "Слайд 13: В результате получен прототип сервиса MindHelper: сайт, Telegram-бот, backend API, PostgreSQL-база, локальная интеграция Qwen3 и safety-flow. Практическая значимость работы состоит в создании основы для сервиса первичной психологической поддержки, который может развиваться за счёт экспертной разметки, расширения red-team корпуса и дальнейшей LoRA-адаптации.",
        "Слайд 14: Спасибо за внимание."
    ]
    doc=Document()
    st=doc.styles['Normal']; st.font.name='Times New Roman'; st.font.size=Pt(12)
    title=doc.add_paragraph('Текст выступления к презентации')
    title.alignment=WD_ALIGN_PARAGRAPH.CENTER
    for r in title.runs: r.bold=True; r.font.name='Times New Roman'; r.font.size=Pt(14)
    for text in paras:
        p=doc.add_paragraph(text); p.alignment=WD_ALIGN_PARAGRAPH.JUSTIFY
        p.paragraph_format.first_line_indent=Pt(28); p.paragraph_format.space_after=Pt(6)
        for r in p.runs: r.font.name='Times New Roman'; r.font.size=Pt(12)
    doc.save(OUT_SCRIPT)


if __name__ == '__main__':
    build_deck(); build_speech()
    print(OUT_PPTX)
    print(OUT_SCRIPT)




